import os
import re
import math
import logging
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.documents import Document

IMAGES_DIR = os.getenv("IMAGES_DIR", "./uploads/images")
os.makedirs(IMAGES_DIR, exist_ok=True)

# Min image size in bytes to skip tiny icons/logos
_MIN_IMAGE_BYTES = int(os.getenv("RAG_MIN_IMAGE_BYTES", "3000"))

logger = logging.getLogger(__name__)

# Configuration for Chroma — override with env CHROMA_PERSIST_DIR if you change embedding model
CHROMA_PERSIST_DIR = os.getenv("CHROMA_PERSIST_DIR", "./chroma_db")

# sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2 :
#   - 457 MB (deja telecharge, chargement ~15s)
#   - 50+ langues dont FR et EN, dim 384
# Pour changer de modele : supprimer chroma_db et reimporter les docs.
EMBEDDING_MODEL_NAME = os.getenv(
    "RAG_EMBEDDING_MODEL",
    "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
)

# normalize pour une similarité cosinus correcte
EMBEDDING_ENCODE_KWARGS = {"normalize_embeddings": True}

SUPPORTED_EXTENSIONS = {".pdf", ".doc", ".docx", ".xlsx", ".xls", ".pptx", ".ppt", ".png", ".jpg", ".jpeg"}

# Lazy singletons (avoid loading heavy models at import time during tooling)
_embeddings = None

# Cross-encoder reranker is DISABLED to prevent high CPU/RAM usage on startup
# that could freeze the PC. Vector search alone provides good results.
# To re-enable: set RAG_ENABLE_RERANKER=1 in environment variables.
_RERANKER_ENABLED = os.getenv("RAG_ENABLE_RERANKER", "0").strip() == "1"
_reranker = None


def get_embeddings():
    global _embeddings
    if _embeddings is None:
        logger.info("Loading embedding model: %s", EMBEDDING_MODEL_NAME)
        _embeddings = HuggingFaceEmbeddings(
            model_name=EMBEDDING_MODEL_NAME,
            model_kwargs={"trust_remote_code": True},
            encode_kwargs=EMBEDDING_ENCODE_KWARGS,
        )
    return _embeddings


def get_reranker():
    """Cross-encoder reranker — disabled by default to avoid PC freeze on startup."""
    if not _RERANKER_ENABLED:
        return False
    global _reranker
    if _reranker is None:
        try:
            from sentence_transformers import CrossEncoder
            logger.info("Loading cross-encoder reranker")
            _reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
        except Exception as e:
            logger.warning("Cross-encoder not available: %s", e)
            _reranker = False
    return _reranker


def get_vector_store():
    return Chroma(persist_directory=CHROMA_PERSIST_DIR, embedding_function=get_embeddings())


def _ce_score_to_distance(score: float) -> float:
    """Map cross-encoder logit to a pseudo-distance (lower = better) for downstream exp(-d)."""
    rel = 1.0 / (1.0 + math.exp(-float(score)))
    return 3.0 * (1.0 - rel)


def _load_excel(file_path: str) -> list[Document]:
    """Load an Excel workbook — each sheet becomes a document."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")
    wb = openpyxl.load_workbook(file_path, data_only=True)
    docs: list[Document] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                rows.append(line)
        if rows:
            docs.append(Document(
                page_content="\n".join(rows),
                metadata={"sheet": sheet_name},
            ))
    return docs


def _load_pptx(file_path: str) -> list[Document]:
    """Load a PowerPoint presentation — each slide becomes a document."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(file_path)
    docs: list[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"slide": i},
            ))
    return docs


def extract_images_from_pdf(file_path: str, doc_id: int, metadata: dict) -> list[dict]:
    """
    Extract images/diagrams from a PDF using PyMuPDF.
    Saves each image to uploads/images/{doc_id}/ and returns metadata list.
    Returns list of {image_path, page, filename, index}.
    """
    extracted = []
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed — skipping image extraction. Run: pip install pymupdf")
        return []

    doc_images_dir = os.path.join(IMAGES_DIR, str(doc_id))
    os.makedirs(doc_images_dir, exist_ok=True)

    try:
        pdf = fitz.open(file_path)
        img_index = 0
        for page_num in range(len(pdf)):
            page = pdf[page_num]
            image_list = page.get_images(full=True)
            for img_num, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = pdf.extract_image(xref)
                    image_bytes = base_image["image"]
                    # Skip tiny images (icons, logos, decorations)
                    if len(image_bytes) < _MIN_IMAGE_BYTES:
                        continue
                    ext = base_image.get("ext", "png")
                    img_filename = f"page_{page_num + 1}_img_{img_num + 1}.{ext}"
                    img_path = os.path.join(doc_images_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    extracted.append({
                        "image_path": img_path,
                        "image_url_path": f"{doc_id}/{img_filename}",
                        "page": page_num + 1,
                        "filename": metadata.get("filename", ""),
                        "doc_id": str(doc_id),
                        "index": img_index,
                    })
                    img_index += 1
                except Exception as e:
                    logger.warning("Image extraction failed xref=%s page=%s: %s", xref, page_num, e)
        pdf.close()
        logger.info("Extracted %d images from %s (doc_id=%s)", len(extracted), file_path, doc_id)
    except Exception as e:
        logger.warning("PDF image extraction failed for %s: %s", file_path, e)

    return extracted


def _ingest_image(file_path: str, doc_id: int, metadata: dict) -> None:
    """
    Index a standalone image file (PNG/JPG) directly into ChromaDB.
    Copies it to uploads/images/{doc_id}/ and creates a searchable text chunk.
    """
    import shutil
    _, ext = os.path.splitext(file_path)
    filename = metadata.get("filename", os.path.basename(file_path))
    doc_images_dir = os.path.join(IMAGES_DIR, str(doc_id))
    os.makedirs(doc_images_dir, exist_ok=True)

    img_filename = f"img_1{ext.lower()}"
    dest_path = os.path.join(doc_images_dir, img_filename)
    shutil.copy2(file_path, dest_path)

    image_url_path = f"{doc_id}/{img_filename}"
    desc = (
        f"[IMAGE] Diagramme ou figure : {filename}. "
        f"Cette image est un diagramme, schema, figure, graphique ou illustration."
    )
    img_meta = dict(metadata)
    img_meta.update({
        "doc_id": str(doc_id),
        "type": "image",
        "image_url_path": image_url_path,
        "page": 1,
    })
    vector_store = get_vector_store()
    vector_store.add_documents([Document(page_content=desc, metadata=img_meta)])
    logger.info("Indexed image file %s as doc_id=%s", filename, doc_id)


def ingest_document(file_path: str, doc_id: int, metadata: dict):
    """
    Extracts text from a document, chunks it, and adds it to the vector store.
    Also extracts images from PDFs and indexes them as searchable chunks.
    Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, PNG, JPG, JPEG.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    # Standalone image files — index directly as image chunks
    if ext in (".png", ".jpg", ".jpeg"):
        _ingest_image(file_path, doc_id, metadata)
        return

    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        documents = loader.load()
    elif ext in [".docx", ".doc"]:
        loader = Docx2txtLoader(file_path)
        documents = loader.load()
    elif ext in [".xlsx", ".xls"]:
        documents = _load_excel(file_path)
    elif ext in [".pptx", ".ppt"]:
        documents = _load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

    for doc in documents:
        doc.metadata.update(metadata)
        doc.metadata["doc_id"] = str(doc_id)

    # QMS-friendly splitting: headings, paragraphs, then sentences.
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=900,
        chunk_overlap=220,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = text_splitter.split_documents(documents)

    vector_store = get_vector_store()

    if chunks:
        vector_store.add_documents(chunks)

    # Extract and index images for PDFs
    if ext == ".pdf":
        images = extract_images_from_pdf(file_path, doc_id, metadata)
        image_chunks = []
        for img in images:
            # Build a descriptive text that will match queries about diagrams/figures
            desc = (
                f"[IMAGE] Page {img['page']} — Diagramme ou figure dans le document "
                f"{img['filename']}. "
                f"Cette image est un diagramme, schema, figure, graphique ou illustration "
                f"extraite de la page {img['page']}."
            )
            img_meta = dict(metadata)
            img_meta.update({
                "doc_id": str(doc_id),
                "type": "image",
                "image_url_path": img["image_url_path"],
                "page": img["page"],
            })
            image_chunks.append(Document(page_content=desc, metadata=img_meta))
        if image_chunks:
            vector_store.add_documents(image_chunks)
            logger.info("Indexed %d image chunks for doc_id=%s", len(image_chunks), doc_id)


class SimpleBM25:
    def __init__(self, corpus: list[list[str]], k1: float = 1.5, b: float = 0.75):
        self.k1 = k1
        self.b = b
        self.corpus_size = len(corpus)
        self.avg_doc_len = sum(len(doc) for doc in corpus) / max(self.corpus_size, 1)
        self.doc_lens = [len(doc) for doc in corpus]
        
        # Document term frequencies
        self.doc_term_freqs = []
        # Document frequency for terms
        self.doc_freqs = {}
        
        for doc in corpus:
            freqs = {}
            for term in doc:
                freqs[term] = freqs.get(term, 0) + 1
            self.doc_term_freqs.append(freqs)
            for term in freqs:
                self.doc_freqs[term] = self.doc_freqs.get(term, 0) + 1
                
    def get_idf(self, term: str) -> float:
        df = self.doc_freqs.get(term, 0)
        return math.log(1.0 + (self.corpus_size - df + 0.5) / (df + 0.5))

    def get_scores(self, query: list[str]) -> list[float]:
        scores = []
        for i in range(self.corpus_size):
            score = 0.0
            doc_len = self.doc_lens[i]
            tf_dict = self.doc_term_freqs[i]
            for term in query:
                if term in tf_dict:
                    tf = tf_dict[term]
                    idf = self.get_idf(term)
                    num = tf * (self.k1 + 1.0)
                    denom = tf + self.k1 * (1.0 - self.b + self.b * (doc_len / self.avg_doc_len))
                    score += idf * (num / denom)
            scores.append(score)
        return scores


def tokenize(text: str) -> list[str]:
    return [w for w in re.findall(r'\w+', (text or "").lower()) if w]


def search_similar_chunks(
    query: str,
    k: int = 4,
    metadata_filter: dict | None = None,
    fetch_multiplier: int = 6,
):
    """
    Hybrid search: Cosine vector similarity + BM25 keyword search blended via Reciprocal Rank Fusion (RRF).
    Then optional cross-encoder rerank for better precision.
    Returns (Document, distance) with lower = better.
    """
    vector_store = get_vector_store()
    n_fetch = min(max(k * fetch_multiplier, k + 8), 48)
    
    # 1. Fetch vector search results
    try:
        vector_hits_with_score = vector_store.similarity_search_with_score(
            query, 
            k=n_fetch, 
            filter=metadata_filter
        )
    except Exception as e:
        logger.warning("Vector search failed: %s", e)
        vector_hits_with_score = []

    # 2. Fetch all documents from collection to run BM25
    col = vector_store._collection
    try:
        if metadata_filter:
            all_data = col.get(where=metadata_filter)
        else:
            all_data = col.get()
    except Exception as e:
        logger.warning("Col.get failed: %s", e)
        all_data = {"documents": [], "metadatas": [], "ids": []}

    documents = all_data.get('documents') or []
    metadatas = all_data.get('metadatas') or []
    ids = all_data.get('ids') or []

    all_docs = []
    for text, meta, doc_id in zip(documents, metadatas, ids):
        all_docs.append(Document(page_content=text, metadata=meta))

    # 3. Calculate BM25 scores
    bm25_docs = []
    if all_docs:
        tokenized_query = tokenize(query)
        tokenized_corpus = [tokenize(doc.page_content) for doc in all_docs]
        bm25 = SimpleBM25(tokenized_corpus)
        bm25_scores = bm25.get_scores(tokenized_query)
        
        bm25_hits = list(zip(all_docs, bm25_scores))
        bm25_hits.sort(key=lambda x: x[1], reverse=True)
        bm25_docs = [doc for doc, score in bm25_hits[:n_fetch] if score > 0.0]

    # 4. Perform Reciprocal Rank Fusion (RRF)
    vector_hits = [doc for doc, _ in vector_hits_with_score]
    
    def get_doc_key(doc: Document) -> str:
        meta = doc.metadata or {}
        return f"{meta.get('doc_id')}_{meta.get('page')}_{doc.page_content[:100]}"
    
    doc_map = {}
    rrf_scores = {}
    
    for rank, doc in enumerate(vector_hits, start=1):
        key = get_doc_key(doc)
        doc_map[key] = doc
        rrf_scores[key] = rrf_scores.get(key, 0.0) + 1.0 / (60.0 + rank)
        
    for rank, doc in enumerate(bm25_docs, start=1):
        key = get_doc_key(doc)
        doc_map[key] = doc
        rrf_scores[key] = rrf_scores.get(key, 0.0) + 1.0 / (60.0 + rank)

    if not rrf_scores:
        return []

    # Sort candidates by RRF score descending
    sorted_keys = sorted(rrf_scores.keys(), key=lambda key: rrf_scores[key], reverse=True)
    top_candidates = [doc_map[key] for key in sorted_keys[:n_fetch]]

    # 5. Map RRF scores to pseudo-distances, blended with actual vector distance if available
    distance_map = {get_doc_key(doc): dist for doc, dist in vector_hits_with_score}
    max_rrf = 2.0 / 61.0
    min_rrf = 1.0 / 108.0
    span = max_rrf - min_rrf
    
    blended_hits = []
    for doc in top_candidates:
        key = get_doc_key(doc)
        rrf_val = rrf_scores[key]
        
        ratio = (max_rrf - rrf_val) / max(span, 1e-5)
        ratio = max(0.0, min(ratio, 1.0))
        pseudo_dist = 0.2 + 1.6 * ratio
        
        if key in distance_map:
            actual_dist = distance_map[key]
            # Blend actual cosine distance (70%) and RRF pseudo-distance (30%)
            final_dist = 0.7 * actual_dist + 0.3 * pseudo_dist
        else:
            # BM25 only chunk
            final_dist = pseudo_dist
            
        blended_hits.append((doc, final_dist))

    # 6. Optional Cross-Encoder reranking
    reranker = get_reranker()
    if reranker is False:
        return blended_hits[:k]

    try:
        max_chars = 1500
        pairs = []
        for doc, _dist in blended_hits:
            text = (doc.page_content or "")[:max_chars]
            pairs.append([query, text])
        scores = reranker.predict(pairs, show_progress_bar=False)
        rescored = []
        for (doc, orig_dist), score in zip(blended_hits, scores):
            pseudo_dist = _ce_score_to_distance(float(score))
            blend = 0.75 * pseudo_dist + 0.25 * float(orig_dist)
            rescored.append((doc, blend))
        rescored.sort(key=lambda x: x[1])
        return rescored[:k]
    except Exception as e:
        logger.warning("Reranking failed, using hybrid RRF order: %s", e)
        return blended_hits[:k]


def remove_document_from_index(doc_id: int):
    vector_store = get_vector_store()
    try:
        vector_store._collection.delete(where={"doc_id": str(doc_id)})
    except Exception as e:
        print(f"Error removing document from index: {e}")
