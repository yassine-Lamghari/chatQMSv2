"""
Script de génération des documents de test QMS.
Génère : PDF, DOCX, XLSX, PPTX avec contenu ISO 9001 réaliste.
"""

import os
import sys

# Ajouter le backend au path pour utiliser les libs installées dans le venv
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DOCS_TEST_DIR = os.path.join(BASE_DIR, "docs_test")
os.makedirs(DOCS_TEST_DIR, exist_ok=True)

# ─── 1. PDF : Procédure Qualité ISO 9001 ──────────────────────────────────────
def generate_pdf():
    from fpdf import FPDF
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "PROCEDURE QUALITE - ISO 9001:2015", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 6, "Reference : PQ-001 | Version : 3.0 | Date : 2024-01-15", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.cell(0, 6, "Proprietaire : Responsable Qualite | Criticite : High", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(8)

    sections = [
        ("1. OBJET ET DOMAINE D'APPLICATION",
         "Cette procedure definit les exigences relatives au systeme de management de la qualite (SMQ) "
         "conformement a la norme ISO 9001:2015. Elle s'applique a l'ensemble des processus de conception, "
         "production, controle et livraison des produits et services de l'organisation. "
         "Tout ecart doit etre documente via une fiche de non-conformite (FNC)."),
        ("2. REFERENCES NORMATIVES",
         "- ISO 9001:2015 -- Systemes de management de la qualite -- Exigences\n"
         "- ISO 9000:2015 -- Vocabulaire et principes essentiels\n"
         "- ISO 19011:2018 -- Lignes directrices pour l'audit des systemes de management\n"
         "- EN ISO 13485:2016 -- Dispositifs medicaux -- SMQ (si applicable)"),
        ("3. RESPONSABILITES",
         "Responsable Qualite (RQ) : Maintenir et ameliorer le SMQ, animer les revues de direction.\n"
         "Directeur General (DG) : Approuver la politique qualite, allouer les ressources.\n"
         "Chefs de service : Appliquer les procedures dans leur perimetre, signaler les NC.\n"
         "Auditeurs internes : Realiser les audits selon le programme annuel valide par le RQ."),
        ("4. DESCRIPTION DU PROCESSUS",
         "4.1 Planification (Clause 6) : Identifier les risques et opportunites, definir les objectifs qualite "
         "SMART (Specifiques, Mesurables, Atteignables, Realistes, Temporels). Les objectifs sont revus "
         "trimestriellement en reunion de pilotage.\n\n"
         "4.2 Support (Clause 7) : Gerer les competences, la sensibilisation, la communication et "
         "l'information documentee. Toute modification d'un document maitre requiert l'approbation du RQ.\n\n"
         "4.3 Realisation (Clause 8) : Maitriser la conception, les achats, la production et la livraison. "
         "Les criteres d'acceptance sont definis dans les fiches de controle associees.\n\n"
         "4.4 Evaluation des performances (Clause 9) : Surveiller, mesurer, analyser et evaluer via "
         "indicateurs KPI. Audit interne minimum une fois par an. Revue de direction annuelle obligatoire.\n\n"
         "4.5 Amelioration (Clause 10) : Traiter les non-conformites, mener des actions correctives, "
         "ameliorer continuellement le SMQ par la roue de Deming (PDCA)."),
        ("5. INDICATEURS QUALITE (KPI)",
         "- Taux de non-conformite produit : objectif < 2% par mois\n"
         "- Satisfaction client (CSAT) : objectif >= 85/100\n"
         "- Delai de traitement des reclamations : objectif <= 5 jours ouvrables\n"
         "- Taux de realisation du plan d'audit : objectif >= 90%\n"
         "- Efficacite des actions correctives (taux de recloture) : objectif >= 80%"),
        ("6. GESTION DES NON-CONFORMITES",
         "Toute non-conformite (NC) detectable doit etre enregistree dans le systeme de gestion des NC "
         "dans les 24h suivant sa detection. Une analyse des causes (methode 5 Pourquoi ou diagramme "
         "d'Ishikawa) est obligatoire pour les NC de criticite High ou Critical. "
         "L'action corrective doit etre definie, mise en oeuvre et verifiee dans un delai maximal de 30 jours. "
         "Les NC recurrentes (>3 fois en 6 mois) declenchent automatiquement une revue de processus."),
        ("7. ENREGISTREMENTS OBLIGATOIRES",
         "- FQE-001 : Fiche de non-conformite\n"
         "- FQE-002 : Plan d'action correctif/preventif (CAPA)\n"
         "- FQE-003 : Rapport d'audit interne\n"
         "- FQE-004 : Compte-rendu de revue de direction\n"
         "- FQE-005 : Registre de formation et habilitation"),
        ("8. HISTORIQUE DES REVISIONS",
         "Version 1.0 (2022-03-01) : Creation initiale.\n"
         "Version 2.0 (2023-06-15) : Ajout section KPI et integration ISO 9001:2015 clause 9.\n"
         "Version 3.0 (2024-01-15) : Mise a jour section 6 (delais NC), ajout FQE-005."),
    ]

    for title, body in sections:
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_fill_color(220, 230, 245)
        pdf.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT", fill=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 6, body)
        pdf.ln(4)

    out = os.path.join(DOCS_TEST_DIR, "procedure_qualite_ISO9001.pdf")
    pdf.output(out)
    print(f"[OK] PDF genere : {out}")


# ─── 2. DOCX : Checklist Audit Qualite ────────────────────────────────────────
def generate_docx():
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = DocxDocument()
    doc.core_properties.author = "Responsable Qualite"
    doc.core_properties.title = "Checklist Audit Qualite Interne ISO 9001:2015"

    h = doc.add_heading("CHECKLIST AUDIT QUALITE INTERNE -- ISO 9001:2015", 0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("Reference : AUDIT-CK-001 | Version : 2.1 | Date : 2024-02-20")
    doc.add_paragraph("Auditeur : ________________________  Site : ________________________")
    doc.add_paragraph("Service audite : __________________  Date audit : ________________")
    doc.add_paragraph("")

    sections = {
        "CONTEXTE DE L'ORGANISATION (Clause 4)": [
            ("4.1", "Les enjeux internes et externes pertinents sont-ils identifies et surveilles ?"),
            ("4.2", "Les parties interessees pertinentes et leurs exigences sont-elles documentees ?"),
            ("4.3", "Le domaine d'application du SMQ est-il defini et maintenu a jour ?"),
            ("4.4", "Les processus necessaires au SMQ sont-ils identifies avec leurs interactions ?"),
        ],
        "LEADERSHIP (Clause 5)": [
            ("5.1", "La direction demontre-t-elle son engagement envers le SMQ ?"),
            ("5.2", "La politique qualite est-elle documentee, communiquee et comprise ?"),
            ("5.3", "Les roles, responsabilites et autorites sont-ils clairement definis ?"),
        ],
        "PLANIFICATION (Clause 6)": [
            ("6.1", "Les risques et opportunites sont-ils determines et traites ?"),
            ("6.2", "Les objectifs qualite sont-ils mesurables et surveilles ?"),
            ("6.3", "Les modifications du SMQ sont-elles planifiees et maitrisees ?"),
        ],
        "SUPPORT (Clause 7)": [
            ("7.1", "Les ressources (humaines, infrastructure, environnement) sont-elles adequates ?"),
            ("7.2", "Les competences du personnel sont-elles determinees et verifiees ?"),
            ("7.3", "Le personnel est-il sensibilise a la politique qualite et aux NC ?"),
            ("7.5", "L'information documentee est-elle creee, mise a jour et maitrisee ?"),
        ],
        "REALISATION (Clause 8)": [
            ("8.1", "La planification et la maitrise operationnelles sont-elles en place ?"),
            ("8.2", "Les exigences relatives aux produits/services sont-elles determinees ?"),
            ("8.4", "La maitrise des processus, produits et services fournis par des prestataires est-elle assuree ?"),
            ("8.5", "La production et la prestation de service sont-elles maitrisees ?"),
            ("8.7", "Les elements de sortie non conformes sont-ils traites ?"),
        ],
        "EVALUATION DES PERFORMANCES (Clause 9)": [
            ("9.1", "Les indicateurs de surveillance et de mesure sont-ils suivis ?"),
            ("9.2", "Des audits internes sont-ils programmes et realises ?"),
            ("9.3", "La revue de direction est-elle realisee et documentee ?"),
        ],
        "AMELIORATION (Clause 10)": [
            ("10.2", "Les non-conformites sont-elles traitees et des actions correctives menees ?"),
            ("10.3", "L'amelioration continue est-elle demontrée par les resultats des indicateurs ?"),
        ],
    }

    for section_title, items in sections.items():
        doc.add_heading(section_title, level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        hdr[0].text = "Ref."
        hdr[1].text = "Critere d'audit"
        hdr[2].text = "Conforme (O/N/NA)"
        hdr[3].text = "Observations / Preuves"
        for ref, question in items:
            row = table.add_row().cells
            row[0].text = ref
            row[1].text = question
            row[2].text = ""
            row[3].text = ""
        doc.add_paragraph("")

    doc.add_heading("BILAN DE L'AUDIT", level=1)
    doc.add_paragraph("Points forts identifies :")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("Non-conformites detectees :")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("Observations et opportunites d'amelioration :")
    doc.add_paragraph("_" * 80)
    doc.add_paragraph("Conclusion (Favorable / Defavorable / Avec reserves) : ____________")

    out = os.path.join(DOCS_TEST_DIR, "checklist_audit_ISO9001.docx")
    doc.save(out)
    print(f"[OK] DOCX genere : {out}")


# ─── 3. XLSX : Plan de Controle Qualite ───────────────────────────────────────
def generate_xlsx():
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()

    # ── Feuille 1 : Plan de contrôle ──
    ws1 = wb.active
    ws1.title = "Plan de Controle"

    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=11)
    thin = Side(style="thin")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    headers = [
        "N° Etape", "Etape de production", "Caracteristique controllee",
        "Methode de controle", "Instrument", "Frequence",
        "Valeur nominale", "Tolerance (+/-)", "Taille echantillon",
        "Critere d'acceptation", "Enregistrement", "Responsable"
    ]

    for col, h in enumerate(headers, 1):
        cell = ws1.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border

    data = [
        [1, "Reception matiere premiere", "Dimensions geometriques", "Mesure directe", "Pied a coulisse digital", "100%", 25.00, 0.05, 5, "Cpk >= 1.33", "FQC-001", "Operateur reception"],
        [2, "Reception matiere premiere", "Durete Shore A", "Durometre", "Durometre Shore A", "Par lot", 65, 3, 3, "60-70 Shore A", "FQC-001", "Labo qualite"],
        [3, "Decoupe", "Longueur piece", "Mesure directe", "Metre ruban etalonné", "10%", 120.0, 0.2, 10, "119.8-120.2 mm", "FQC-002", "Operateur decoupe"],
        [4, "Decoupe", "Angle de coupe", "Mesure angulaire", "Rapporteur d'angle", "5%", 90.0, 0.5, 5, "89.5-90.5 deg", "FQC-002", "Controleur Q"],
        [5, "Assemblage", "Force d'insertion", "Essai fonctionnel", "Dynamometre", "100%", 50.0, 5.0, 1, "45-55 N", "FQC-003", "Operateur montage"],
        [6, "Assemblage", "Etancheite", "Test helium", "Detecteur fuites He", "100%", 0, "< 1e-6 mbar.l/s", 1, "Pass/Fail", "FQC-003", "Controleur Q"],
        [7, "Traitement thermique", "Temperature four", "Thermocouple calibre", "Enregistreur T°", "Continu", 180.0, 5.0, "Continu", "175-185 degC", "FQC-004", "Operateur four"],
        [8, "Traitement thermique", "Duree cycle", "Chronometre", "Minuteur certifie", "100%", 45, 2, 1, "43-47 min", "FQC-004", "Operateur four"],
        [9, "Controle final", "Aspect visuel", "Inspection visuelle", "Eclairage 1000 lux", "100%", "Conforme", "Cf. gamme visuelle", 1, "0 defaut majeur", "FQC-005", "Controleur final"],
        [10, "Controle final", "Marquage etiquette", "Lecture code barres", "Scanner Datalogic", "100%", "Conforme", "N/A", 1, "Lecture OK + traçabilite", "FQC-005", "Exp. produit fini"],
        [11, "Emballage", "Resistance emballage", "Test chute 1m", "Table + regle", "Par lot", "Integre", "N/A", 3, "0 dommage", "FQC-006", "Responsable exp."],
    ]

    for row_idx, row_data in enumerate(data, 2):
        fill_color = "EBF3FB" if row_idx % 2 == 0 else "FFFFFF"
        row_fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
        for col_idx, val in enumerate(row_data, 1):
            cell = ws1.cell(row=row_idx, column=col_idx, value=val)
            cell.fill = row_fill
            cell.border = border
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    col_widths = [8, 22, 26, 22, 22, 12, 14, 14, 12, 22, 12, 20]
    for i, w in enumerate(col_widths, 1):
        ws1.column_dimensions[get_column_letter(i)].width = w
    ws1.row_dimensions[1].height = 35

    # ── Feuille 2 : KPI Qualite ──
    ws2 = wb.create_sheet("KPI Qualite")
    ws2.cell(row=1, column=1, value="TABLEAU DE BORD KPI QUALITE -- 2024").font = Font(bold=True, size=14)
    kpi_headers = ["KPI", "Objectif", "Jan", "Fev", "Mar", "Avr", "Mai", "Jun", "Jul", "Aou", "Sep", "Oct", "Nov", "Dec"]
    for col, h in enumerate(kpi_headers, 1):
        c = ws2.cell(row=3, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
        c.alignment = Alignment(horizontal="center")
    kpi_data = [
        ["Taux NC produit (%)", "< 2%", 1.8, 1.5, 2.1, 1.9, 1.7, 1.4, 1.6, 1.3, 1.8, 2.0, 1.5, 1.7],
        ["CSAT client (/100)", ">= 85", 87, 86, 84, 88, 90, 89, 87, 91, 88, 86, 90, 92],
        ["Delai traitement NC (j)", "<= 5j", 4.2, 5.1, 3.8, 4.5, 3.9, 4.1, 5.0, 3.6, 4.8, 4.3, 3.7, 4.0],
        ["Taux realisation audits (%)", ">= 90%", 100, 100, 83, 100, 100, 83, 100, 100, 100, 100, 83, 100],
        ["Efficacite CAPA (%)", ">= 80%", 82, 85, 78, 88, 84, 91, 86, 89, 83, 87, 90, 88],
    ]
    for row_idx, row in enumerate(kpi_data, 4):
        for col_idx, val in enumerate(row, 1):
            ws2.cell(row=row_idx, column=col_idx, value=val)

    out = os.path.join(DOCS_TEST_DIR, "plan_controle_qualite.xlsx")
    wb.save(out)
    print(f"[OK] XLSX genere : {out}")


# ─── 4. PPTX : Formation Qualite ──────────────────────────────────────────────
def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(layout_idx, title_text, body_lines=None, subtitle=None):
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                if body_lines:
                    tf = ph.text_frame
                    tf.text = body_lines[0]
                    tf.paragraphs[0].font.size = Pt(18)
                    for line in body_lines[1:]:
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(16)
                elif subtitle:
                    ph.text = subtitle
        return slide

    # Slide 1 : Titre
    s = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(s)
    slide1.shapes.title.text = "FORMATION QUALITE ISO 9001:2015"
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Sensibilisation du personnel -- Module QMS-F001\nDuree : 4 heures | Niveau : Tout le personnel"

    slides_data = [
        (1, "OBJECTIFS DE LA FORMATION", [
            "A l'issue de cette formation, le participant sera capable de :",
            "  • Expliquer les principes de la norme ISO 9001:2015",
            "  • Identifier sa contribution aux processus qualite",
            "  • Signaler et documenter une non-conformite",
            "  • Comprendre les consequences d'un ecart qualite",
            "  • Appliquer la demarche PDCA dans son activite quotidienne",
        ]),
        (1, "LES 7 PRINCIPES DU MANAGEMENT QUALITE", [
            "1. Orientation client -- Satisfaire et depasser les attentes",
            "2. Leadership -- La direction cree les conditions de reussite",
            "3. Implication du personnel -- Le capital humain est cle",
            "4. Approche processus -- Penser flux et interactions",
            "5. Amelioration -- PDCA permanent (Plan-Do-Check-Act)",
            "6. Prise de decision fondee sur preuves -- Donnees, indicateurs, faits",
            "7. Management des relations -- Partenaires et fournisseurs inclus",
        ]),
        (1, "IDENTIFICATION ET GESTION DES NON-CONFORMITES", [
            "Definition : Tout ecart par rapport aux exigences specifiees",
            "",
            "Processus de traitement (delai max 30 jours) :",
            "  1. Detection et signalement immediat (< 24h)",
            "  2. Enregistrement dans le systeme NC (FQE-001)",
            "  3. Isolation/quarantaine du produit non-conforme",
            "  4. Analyse des causes (5 Pourquoi / Ishikawa)",
            "  5. Action corrective (CAPA) -- FQE-002",
            "  6. Verification d'efficacite et cloture",
        ]),
        (1, "LA ROUE DE DEMING -- PDCA", [
            "PLAN (Planifier) -- Definir objectifs et methodes",
            "  Ex : Reduire le taux de NC de 2% a 1.5% en 6 mois",
            "",
            "DO (Realiser) -- Mettre en oeuvre le plan d'action",
            "  Ex : Former les operateurs, modifier la gamme de controle",
            "",
            "CHECK (Verifier) -- Mesurer les resultats",
            "  Ex : Analyser le taux NC mensuel, comparer a l'objectif",
            "",
            "ACT (Agir) -- Standardiser si OK, corriger si NON",
            "  Ex : Mettre a jour la procedure, partager les bonnes pratiques",
        ]),
        (1, "KPI QUALITE -- INDICATEURS CLES", [
            "Taux de non-conformite produit : objectif < 2% / mois",
            "Satisfaction client (CSAT) : objectif >= 85/100",
            "Delai traitement reclamations : objectif <= 5 jours ouvrables",
            "Taux realisation plan d'audit : objectif >= 90%",
            "Efficacite des CAPA (taux de recloture) : objectif >= 80%",
            "",
            "Suivi mensuel en reunion de pilotage qualite.",
            "Revue annuelle en Revue de Direction (obligatoire ISO 9001).",
        ]),
        (1, "EVALUATION ET CERTIFICATION", [
            "QCM de 20 questions -- Seuil de validation : 14/20 (70%)",
            "",
            "Themes evalues :",
            "  • Principes ISO 9001 (5 questions)",
            "  • Gestion des non-conformites (5 questions)",
            "  • PDCA et amelioration continue (5 questions)",
            "  • KPI et indicateurs (3 questions)",
            "  • Cas pratiques (2 questions)",
            "",
            "Attestation de formation delivree si score >= 70%",
            "Recyclage obligatoire tous les 2 ans (FQE-005)",
        ]),
    ]

    for layout_idx, title, body in slides_data:
        add_slide(layout_idx, title, body)

    out = os.path.join(DOCS_TEST_DIR, "formation_qualite_ISO9001.pptx")
    prs.save(out)
    print(f"[OK] PPTX genere : {out}")


# ─── 5. Copie image test ──────────────────────────────────────────────────────
def copy_test_image():
    import shutil
    src = os.path.join(BASE_DIR, "test_assets", "test_diagram.png")
    dst = os.path.join(DOCS_TEST_DIR, "schema_processus_qualite.png")
    if os.path.exists(src):
        shutil.copy2(src, dst)
        print(f"[OK] Image copiee : {dst}")
    else:
        print("[WARN] test_assets/test_diagram.png introuvable -- image ignoree")


# ─── 6. README corpus ─────────────────────────────────────────────────────────
def generate_readme_corpus():
    content = """# Corpus de Documents de Test -- QMS Chatbot RAG

Ce dossier contient les documents utilisés pour tester et valider le système RAG.

## Formats testés

| Fichier | Format | Contenu | Criticité |
|---|---|---|---|
| `procedure_qualite_ISO9001.pdf` | PDF texte | Procédure qualité complète ISO 9001:2015 (8 sections) | High |
| `checklist_audit_ISO9001.docx` | Word (DOCX) | Checklist d'audit interne par clause ISO 9001 | Medium |
| `plan_controle_qualite.xlsx` | Excel (XLSX) | Plan de contrôle (11 étapes) + tableau KPI 2024 | High |
| `formation_qualite_ISO9001.pptx` | PowerPoint (PPTX) | Support de formation (6 slides) | Low |
| `schema_processus_qualite.png` | Image (PNG) | Schéma de processus qualité | Low |

## Utilisation

Ces documents doivent être importés dans l'application via l'interface d'administration :
1. Connectez-vous avec un compte **admin**
2. Allez dans **Gestion des documents**
3. Importez chaque fichier avec les métadonnées appropriées
4. Vérifiez l'indexation dans ChromaDB

## Couverture des formats

| Format | Support annoncé | Document de test fourni | Testé |
|---|---|---|---|
| PDF (texte natif) | ✅ | `procedure_qualite_ISO9001.pdf` | ✅ |
| PDF (scanné/OCR) | ⚠️ Via PyMuPDF (images extraites) | -- | Partiel |
| Word (.docx) | ✅ | `checklist_audit_ISO9001.docx` | ✅ |
| Excel (.xlsx) | ✅ | `plan_controle_qualite.xlsx` | ✅ |
| PowerPoint (.pptx) | ✅ | `formation_qualite_ISO9001.pptx` | ✅ |
| Image (.png/.jpg) | ✅ | `schema_processus_qualite.png` | ✅ |

> **Note** : Les PDF scannés sont traités via l'extraction d'images PyMuPDF.
> Le contenu textuel des PDF scannés n'est pas extrait par OCR dans la version actuelle.
"""
    out = os.path.join(DOCS_TEST_DIR, "README_corpus.md")
    with open(out, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"[OK] README corpus genere : {out}")


if __name__ == "__main__":
    print("=" * 60)
    print("Generation des documents de test QMS")
    print("=" * 60)
    try:
        generate_pdf()
    except Exception as e:
        print(f"[ERREUR] PDF : {e}")
    try:
        generate_docx()
    except Exception as e:
        print(f"[ERREUR] DOCX : {e}")
    try:
        generate_xlsx()
    except Exception as e:
        print(f"[ERREUR] XLSX : {e}")
    try:
        generate_pptx()
    except Exception as e:
        print(f"[ERREUR] PPTX : {e}")
    copy_test_image()
    generate_readme_corpus()
    print("=" * 60)
    print("Termine ! Documents dans : docs_test/")
    print("=" * 60)
